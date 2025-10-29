import streamlit as st
import google.generativeai as genai
import os
import io
import re
from pptx import Presentation
import fitz  # PyMuPDF
from fpdf import FPDF
from PIL import Image

# ------------------------------------------------------------------
# 1. THE SYSTEM PROMPT FOR GEMINI
# ------------------------------------------------------------------
SYSTEM_PROMPT = """
You are an expert academic author and textbook editor. Your sole task is to 
convert the provided raw lecture slide content into a single, comprehensive, 
and formally written textbook chapter.

**Core Directives:**

1.  **ROLE:** Act as an author writing a chapter, not a student taking notes.
2.  **TONE:** Maintain a formal, academic, and descriptive tone. The language must 
    be clear, objective, and educational.
3.  **NO SUMMARIZATION:** You are forbidden from summarizing. Your primary goal is 
    high-fidelity conversion. You must preserve ALL information, data points, 
    and nuances from the original text.
4.  **CONVERT, DON'T LIST:** Transform all bullet points, fragmented phrases, 
    and talking points into complete, well-structured paragraphs. Lists should 
    only be used if the original content is a sequential list of steps.
5.  **ELABORATE CONTEXTUALLY:** When converting fragments (e.g., "• Faster processing"), 
    you must rewrite them as full sentences by inferring context 
    (e.g., "This new algorithm provides significantly faster processing speeds...").

**Formatting Rules:**

1.  **STRUCTURE:**
    * The overall topic shall be the main chapter title (e.g., `# Chapter 7: ...`).
    * Main slide titles shall be treated as major section headings (e.g., `## 7.1 ...`).
    * Key concepts on a slide may be used as subsections (e.g., `### 7.1.1 ...`).

2.  **IMAGE HANDLING:**
    * You will receive image placeholders in the format `[IMAGE: image_filename.png]`.
    * You MUST include these placeholders exactly as-is on their own line.
    * Immediately following the placeholder, you MUST generate a descriptive 
        caption in the textbook style (e.g., *Figure 7.1: A diagram...*).

3.  **KEY TERMS:**
    * When a new key term or acronym is introduced, **bold** it and provide 
        a brief, integrated definition.

**Output Format:**
The final output must be a single, cohesive document formatted in **Markdown**.
"""

# ------------------------------------------------------------------
# 2. HELPER FUNCTIONS (FILE PARSING)
# ------------------------------------------------------------------

def extract_content_from_pptx(file_obj):
    """
    Extracts text and images from a .pptx file.
    """
    st.write("Extracting from PowerPoint...")
    prs = Presentation(file_obj)
    full_text_content = []
    image_dict = {}
    
    for i, slide in enumerate(prs.slides):
        full_text_content.append(f"\n[SLIDE {i+1}]\n")
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text_frame.text)
            
            if shape.shape_type == 13:  # 13 is 'Picture'
                try:
                    image = shape.image
                    img_bytes = image.blob
                    img_ext = image.ext
                    img_filename = f"slide_{i+1}_img_{len(image_dict)}.{img_ext}"
                    
                    image_dict[img_filename] = img_bytes
                    slide_text.append(f"[IMAGE: {img_filename}]")
                except Exception as e:
                    st.warning(f"Could not extract an image from slide {i+1}: {e}")
                    
        full_text_content.append("\n".join(slide_text))
        
    st.write("...PowerPoint extraction complete.")
    return "\n".join(full_text_content), image_dict

def extract_content_from_pdf(file_obj):
    """
    Extracts text and images from a .pdf file.
    """
    st.write("Extracting from PDF...")
    doc = fitz.open(stream=file_obj.read(), filetype="pdf")
    full_text_content = []
    image_dict = {}
    
    for page_num in range(len(doc)):
        full_text_content.append(f"\n[PAGE {page_num+1}]\n")
        page = doc.load_page(page_num)
        
        # Extract text
        full_text_content.append(page.get_text("text"))
        
        # Extract images
        img_list = page.get_images(full=True)
        for img_index, img in enumerate(img_list):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                img_ext = base_image["ext"]
                img_filename = f"page_{page_num+1}_img_{img_index}.{img_ext}"
                
                image_dict[img_filename] = img_bytes
                full_text_content.append(f"[IMAGE: {img_filename}]")
            except Exception as e:
                st.warning(f"Could not extract image {xref} from page {page_num+1}: {e}")

    st.write("...PDF extraction complete.")
    return "\n".join(full_text_content), image_dict

# ------------------------------------------------------------------
# 3. HELPER FUNCTIONS (API & PDF)
# ------------------------------------------------------------------

@st.cache_data
def get_available_models():
    """
    Fetches and filters the list of available Gemini models.
    Returns a list of model names that support 'generateContent'.
    """
    try:
        genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
        models_list = []
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                models_list.append(m.name)
        return models_list
    except Exception as e:
        st.error(f"Error fetching model list: {e}")
        return []

def call_gemini_api(content_to_process, model_name):
    """
    Sends the extracted content and system prompt to the Gemini API
    using the user-selected model.
    """
    try:
        genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
        
        model = genai.GenerativeModel(
            model_name=model_name,
            system_instruction=SYSTEM_PROMPT
        )
        response = model.generate_content(content_to_process)
        return response.text
    except Exception as e:
        st.error(f"Error calling Gemini API: {e}")
        st.stop()


# --- UPDATED PDF CLASS ---
class PDF(FPDF):
    """Custom FPDF class to handle multi-cell with bold/italic"""
    def write_html(self, text):
        # Sanitize text for Unicode
        text = text.replace("’", "'").replace("‘", "'") \
                   .replace("“", '"').replace("”", '"') \
                   .replace("–", "-").replace("—", "-") \
                   .replace("→", "->") \
                   .replace("∪", "U")  # <-- FIX
        
        parts = re.split(r'(\*\*.+?\*\*|_.+?_)', text)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                self.set_font(self.font_family, 'B', self.font_size)
                self.write(self.h, part[2:-2])
            elif part.startswith('_') and part.endswith('_'):
                self.set_font(self.font_family, 'I', self.font_size)
                self.write(self.h, part[1:-1])
            else:
                self.set_font(self.font_family, '', self.font_size)
                self.write(self.h, part)
        self.ln(self.h)


# --- UPDATED PDF CREATION FUNCTION ---
def create_pdf_from_markdown(markdown_text, image_dict):
    """
    Generates a PDF from the LLM's Markdown output, embedding images.
    """
    pdf = PDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_left_margin(15)
    pdf.set_right_margin(15)
    pdf.set_font("Arial", size=12)
    
    lines = markdown_text.split('\n')
    
    for line in lines:
        if line.strip() == "":
            pdf.ln(5)
            continue

        # Sanitize all lines for common Unicode issues
        line = line.replace("’", "'").replace("‘", "'") \
                   .replace("“", '"').replace("”", '"') \
                   .replace("–", "-").replace("—", "-") \
                   .replace("→", "->") \
                   .replace("∪", "U")  # <-- FIX

        # H1 Title: #
        if re.match(r'^#\s', line):
            pdf.set_font("Arial", 'B', 24)
            pdf.multi_cell(0, 12, line[2:].strip(), ln=1)
            pdf.ln(5)
        
        # H2 Title: ##
        elif re.match(r'^##\s', line):
            pdf.set_font("Arial", 'B', 18)
            pdf.multi_cell(0, 10, line[3:].strip(), ln=1)
            pdf.ln(4)

        # H3 Title: ###
        elif re.match(r'^###\s', line):
            pdf.set_font("Arial", 'B', 14)
            pdf.multi_cell(0, 8, line[4:].strip(), ln=1)
            pdf.ln(3)

        # Image Placeholder: [IMAGE: ...]
        elif re.match(r'^\[IMAGE: (.*)\]', line):
            img_filename = re.match(r'^\[IMAGE: (.*)\]', line).group(1).strip()
            if img_filename in image_dict:
                try:
                    img_bytes = image_dict[img_filename]
                    img_file_obj = io.BytesIO(img_bytes)
                    
                    with Image.open(img_file_obj) as img:
                        img_format = img.format.upper()
                        if img_format not in ['JPEG', 'PNG', 'GIF']:
                            with io.BytesIO() as output_bytes:
                                img.save(output_bytes, format="PNG")
                                img_bytes = output_bytes.getvalue()
                                img_format = "PNG"
                        
                        page_width = pdf.w - pdf.l_margin - pdf.r_margin
                        
                        if img.width == 0:
                            continue
                            
                        ratio = img.height / img.width
                        img_width = page_width
                        img_height = page_width * ratio

                        img_file_obj_final = io.BytesIO(img_bytes)
                        pdf.image(img_file_obj_final, x=pdf.get_x(), w=img_width, type=img_format)
                        pdf.ln(img_height + 2) 

                except Exception as e:
                    error_msg = f"[Error embedding image: {img_filename}. {e}]"
                    error_msg = error_msg.replace("’", "'").replace("‘", "'").replace("→", "->").replace("∪", "U")
                    pdf.set_font("Arial", 'I', 10)
                    pdf.set_text_color(255, 0, 0) # Red
                    pdf.multi_cell(0, 5, error_msg, ln=1)
                    pdf.set_text_color(0, 0, 0) # Back to black
            else:
                error_msg = f"[Image not found: {img_filename}]"
                error_msg = error_msg.replace("’", "'").replace("‘", "'").replace("→", "->").replace("∪", "U")
                pdf.set_font("Arial", 'I', 10)
                pdf.set_text_color(255, 0, 0)
                pdf.multi_cell(0, 5, error_msg, ln=1)
                pdf.set_text_color(0, 0, 0)

        # Figure Caption: *Figure...
        elif re.match(r'^\*Figure.*', line):
            pdf.set_font("Arial", 'I', 10)
            pdf.multi_cell(0, 5, line.strip(), ln=1)
            pdf.ln(5)

        # Bullet points
        elif re.match(r'^\s*[\*\-]\s', line):
            pdf.set_font("Arial", '', 12)
            pdf.multi_cell(0, 5, f"  -  {line.lstrip(' *-')}", ln=1) 

        # Numbered lists
        elif re.match(r'^\s*\d+\.\s', line):
            pdf.set_font("Arial", '', 12)
            pdf.multi_cell(0, 5, f"  {line.lstrip()}", ln=1)

        # Paragraph text (with bold/italic)
        else:
            pdf.set_font("Arial", '', 12)
            pdf.write_html(line.strip()) # This function now sanitizes internally
            
    return pdf.output(dest='S')


# ------------------------------------------------------------------
# 4. STREAMLIT APPLICATION UI
# ------------------------------------------------------------------

st.set_page_config(page_title="Lecture to Textbook", layout="wide")
st.title("📚 Lecture to Textbook Generator")
st.markdown("Upload your lecture slides (`.pptx` or `.pdf`) and get a textbook-style PDF back, powered by Gemini.")

# Check for API Key
if "GOOGLE_API_KEY" not in st.secrets:
    st.error("🚨 GOOGLE_API_KEY not found in Streamlit secrets!")
    st.markdown("Please add your API key to a file named `.streamlit/secrets.toml`")
    st.code("GOOGLE_API_KEY = \"YOUR_API_KEY_HERE\"")
    st.stop()

# --- FETCH MODELS AND CREATE DROPDOWN ---
available_models = get_available_models()

if not available_models:
    st.error("Could not load any models. Please check your API key and permissions.")
    st.stop()

def format_model_name(full_name):
    """Helper function to make model names prettier in the dropdown."""
    return full_name.split('/')[-1]

selected_model_name = st.selectbox(
    "Choose your Gemini Model:",
    options=available_models,
    format_func=format_model_name,
    help="Models like '1.5-flash' are faster, '1.5-pro' is more powerful."
)


# File Uploader
uploaded_file = st.file_uploader("Upload your file", type=["pptx", "pdf"])

if uploaded_file is not None:
    
    # Process Button
    if st.button(f"Generate Textbook from {uploaded_file.name}"):
        
        # 1. Extract Content
        with st.spinner(f"Reading and extracting content from {uploaded_file.name}..."):
            try:
                if uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    extracted_text, images = extract_content_from_pptx(uploaded_file)
                elif uploaded_file.type == "application/pdf":
                    extracted_text, images = extract_content_from_pdf(uploaded_file)
                else:
                    st.error("Unsupported file type.")
                    st.stop()
                
                st.success("File extraction complete!")
                st.write(f"Found {len(extracted_text.split())} words and {len(images)} images.")

            except Exception as e:
                st.error(f"Error during file extraction: {e}")
                st.stop()

        # 2. Call Gemini API
        with st.spinner(f"🤖 Calling {format_model_name(selected_model_name)} to write your textbook..."):
            try:
                markdown_output = call_gemini_api(extracted_text, selected_model_name)
                st.success("Gemini has finished writing!")
                
                with st.expander("Show Gemini's Markdown Output"):
                    st.markdown(markdown_output)

            except Exception as e:
                st.error(f"Error calling Gemini: {e}")
                st.stop()
        
        # 3. Create PDF
        with st.spinner("🎨 Formatting your PDF textbook..."):
            try:
                pdf_bytes = create_pdf_from_markdown(markdown_text, images)
                st.success("PDF created successfully!")

                # 4. Download Button
                st.download_button(
                    label="⬇️ Download Your Textbook PDF",
                    data=pdf_bytes,
                    file_name=f"textbook_{uploaded_file.name}.pdf",
                    mime="application/pdf"
                )
            except Exception as e:
                st.error(f"Error creating PDF: {e}")
                st.exception(e) # Show full traceback
                st.stop()
